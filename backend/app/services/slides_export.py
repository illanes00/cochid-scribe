"""Professional PPTX export using python-pptx."""

import re
import struct
from html import unescape
from io import BytesIO
from pathlib import Path
from urllib.parse import unquote, urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Espacio Público palette
PRIMARY_COLOR = RGBColor(0x1A, 0x36, 0x5D)
SECONDARY_COLOR = RGBColor(0xC5, 0x30, 0x30)
ACCENT_COLOR = RGBColor(0x2B, 0x6C, 0xB0)
TEXT_COLOR = RGBColor(0x1A, 0x20, 0x2C)
MUTED_COLOR = RGBColor(0x80, 0x80, 0x80)


def create_presentation(slides_data: dict) -> BytesIO:
    """Create a branded PPTX presentation from slides data."""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = slides_data.get("slides", [])
    theme = slides_data.get("theme", {})

    for slide in slides:
        layout = slide.get("layout", "content")
        if layout == "title":
            add_title_slide(prs, slide, theme)
        else:
            add_content_slide(prs, slide, theme)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def add_title_slide(prs: Presentation, slide: dict, theme: dict):
    """Add a title slide with full-bleed primary color."""

    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    font_name = get_theme_font(theme)
    logo_path = resolve_local_image_path(str(theme.get("logoUrl") or ""))

    # Background block
    shape = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = get_theme_color(theme, "primaryColor", PRIMARY_COLOR)
    shape.line.fill.background()

    title_box = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    p = title_box.text_frame.paragraphs[0]
    p.text = slide.get("title", "")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    subtitle_raw = slide.get("content") or ""
    subtitle_lines = prepare_lines(subtitle_raw)
    subtitle = "\n".join(line for line in subtitle_lines if line.strip())
    if subtitle:
        sub_box = s.shapes.add_textbox(Inches(2), Inches(4.1), Inches(9.333), Inches(1.2))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(20)
        sp.font.name = font_name
        sp.font.color.rgb = RGBColor(235, 235, 235)
        sp.alignment = PP_ALIGN.CENTER

    set_slide_notes(s, slide.get("notes") or "")

    if logo_path:
        add_logo(s, logo_path, left=Inches(10.8), top=Inches(0.6), height=Inches(0.9))


def add_content_slide(prs: Presentation, slide: dict, theme: dict):
    """Add a content slide with header band and body text."""

    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    font_name = get_theme_font(theme)
    logo_path = resolve_local_image_path(str(theme.get("logoUrl") or ""))

    header = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = get_theme_color(theme, "primaryColor", PRIMARY_COLOR)
    header.line.fill.background()

    title_box = s.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = slide.get("title", "")
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.name = font_name
    tp.font.color.rgb = RGBColor(255, 255, 255)

    raw_content = slide.get("content") or ""
    layout = slide.get("layout", "content")
    image_paths = extract_local_image_paths(raw_content)

    body_left = 0.6
    body_top = 1.3
    body_width = 12.1
    body_height = 5.5

    content_left = body_left
    content_top = body_top
    content_width = body_width
    content_height = body_height

    image_left = None
    image_top = None
    image_width = None
    image_height = None

    lines = prepare_lines(raw_content)
    if image_paths and not lines:
        layout = "image-full"

    if image_paths and layout == "image-full":
        image_boxes = get_image_full_boxes(
            body_left, body_top, body_width, body_height, len(image_paths)
        )
        for image_path, box in zip(image_paths, image_boxes, strict=False):
            left, top, width, height = box
            add_picture_fitted(s, image_path, left, top, width, height)
    elif image_paths:
        gutter = 0.3
        side_width = 4.5
        content_width = max(5.0, body_width - side_width - gutter)
        image_left = body_left + content_width + gutter
        image_top = body_top
        image_width = side_width
        image_height = body_height

        add_side_images(
            s,
            image_paths,
            left=image_left,
            top=image_top,
            width=image_width,
            height=image_height,
        )

    if lines:
        content_box = s.shapes.add_textbox(
            Inches(content_left), Inches(content_top), Inches(content_width), Inches(content_height)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        first = True
        for line in lines:
            text = line.rstrip()
            if not text:
                continue
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False

            bullet_match = re.match(r"^(\s*)([-*])\s+(.*)$", text)
            numbered_match = re.match(r"^(\s*)\d+\.\s+(.*)$", text)
            if bullet_match:
                indent = len(bullet_match.group(1))
                para.text = f"• {bullet_match.group(3).strip()}"
                para.level = min(indent // 2, 3)
            elif numbered_match:
                indent = len(numbered_match.group(1))
                para.text = numbered_match.group(2).strip()
                para.level = min(indent // 2, 3)
            else:
                para.text = text.strip()
                para.level = 0

            para.font.size = Pt(18)
            para.font.name = font_name
            para.font.color.rgb = TEXT_COLOR

    footer_box = s.shapes.add_textbox(Inches(0.6), Inches(7), Inches(6), Inches(0.4))
    fp = footer_box.text_frame.paragraphs[0]
    fp.text = "Espacio Público"
    fp.font.size = Pt(10)
    fp.font.name = font_name
    fp.font.color.rgb = MUTED_COLOR

    num_box = s.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
    np = num_box.text_frame.paragraphs[0]
    np.text = str(slide.get("slideNumber", ""))
    np.font.size = Pt(10)
    np.font.name = font_name
    np.font.color.rgb = MUTED_COLOR
    np.alignment = PP_ALIGN.RIGHT

    set_slide_notes(s, slide.get("notes") or "")

    if logo_path:
        add_logo(s, logo_path, left=Inches(11.4), top=Inches(0.1), height=Inches(0.7))


def _parse_hex_color(value: str) -> RGBColor | None:
    value = value.strip()
    if not value.startswith("#"):
        return None
    if len(value) not in {7, 9}:
        return None
    try:
        return RGBColor.from_string(value[1:7])
    except ValueError:
        return None


def get_theme_color(theme: dict, key: str, fallback: RGBColor) -> RGBColor:
    """Resolve a theme color (hex string) by key or return fallback."""
    if not isinstance(theme, dict):
        return fallback
    value = theme.get(key)
    if isinstance(value, str):
        parsed = _parse_hex_color(value)
        if parsed:
            return parsed
    return fallback


def get_theme_font(theme: dict) -> str:
    """Pick the first font from theme fontFamily."""
    if not isinstance(theme, dict):
        return "Arial"
    font = theme.get("fontFamily")
    if not isinstance(font, str) or not font.strip():
        return "Arial"
    return font.split(",")[0].strip().strip("'\"") or "Arial"


def prepare_lines(content: str) -> list[str]:
    """Normalize content and split into lines, preserving bullet intent."""
    normalized = re.sub(r"\r\n?", "\n", content or "")
    stripped = normalized.strip()
    if "<" in stripped and ">" in stripped and re.search(r"<[a-zA-Z][^>]*>", stripped):
        stripped = html_to_markdownish(stripped)
    stripped = re.sub(r"\n{3,}", "\n\n", stripped).strip()
    return stripped.split("\n") if stripped else []


def html_to_markdownish(html: str) -> str:
    """Convert a simple HTML fragment into markdown-like text for slide export."""
    text = html
    # Line breaks
    text = re.sub(r"(?is)<br\\s*/?>", "\n", text)
    text = re.sub(r"(?is)</p\\s*>", "\n", text)
    text = re.sub(r"(?is)</div\\s*>", "\n", text)
    text = re.sub(r"(?is)</h[1-6]\\s*>", "\n", text)

    # Bullets
    text = re.sub(r"(?is)<li[^>]*>\\s*", "- ", text)
    text = re.sub(r"(?is)</li\\s*>", "\n", text)

    # Remove opening container tags we handled via closing tags
    text = re.sub(r"(?is)<p[^>]*>", "", text)
    text = re.sub(r"(?is)<div[^>]*>", "", text)
    text = re.sub(r"(?is)<h[1-6][^>]*>", "", text)
    text = re.sub(r"(?is)<ul[^>]*>", "", text)
    text = re.sub(r"(?is)</ul\\s*>", "\n", text)
    text = re.sub(r"(?is)<ol[^>]*>", "", text)
    text = re.sub(r"(?is)</ol\\s*>", "\n", text)

    # Strip remaining tags
    text = re.sub(r"(?is)<[^>]+>", "", text)

    # Decode entities and normalize whitespace
    text = unescape(text)
    text = re.sub(r"[ \\t]+", " ", text)
    return text.strip()


def export_slides_to_pptx(slides_data: dict, output_path: Path, title: str | None = None) -> None:
    """Save slides_data to a PPTX file on disk."""
    if title and slides_data is not None and isinstance(slides_data, dict):
        slides_data.setdefault("title", title)

    buffer = create_presentation(slides_data or {})
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(buffer.getvalue())


def extract_local_image_paths(html: str, *, limit: int = 8) -> list[Path]:
    """Return local image paths referenced by <img src="...">, preserving order."""
    paths: list[Path] = []
    seen: set[str] = set()
    for src in extract_image_sources(html or ""):
        candidate = resolve_local_image_path(src)
        if not candidate:
            continue
        key = str(candidate)
        if key in seen:
            continue
        seen.add(key)
        paths.append(candidate)
        if len(paths) >= limit:
            break
    return paths


def extract_image_sources(html: str) -> list[str]:
    if not html:
        return []
    return re.findall(r'(?is)<img[^>]+src=["\']([^"\']+)["\']', html)


def resolve_local_image_path(src: str) -> Path | None:
    if not src or src.startswith("data:"):
        return None

    parsed = urlparse(src)
    if parsed.scheme in {"http", "https"}:
        return None

    path_part = unquote(parsed.path or "")
    if not path_part:
        return None

    backend_dir = Path(__file__).resolve().parents[2]
    candidate = (backend_dir / path_part.lstrip("/")).resolve()

    try:
        candidate.relative_to(backend_dir)
    except ValueError:
        return None

    if candidate.is_file():
        return candidate
    return None


def add_logo(slide, logo_path: Path, *, left: Inches, top: Inches, height: Inches) -> None:
    """Add logo image if available."""
    try:
        slide.shapes.add_picture(str(logo_path), left, top, height=height)
    except Exception:
        return


def get_image_full_boxes(
    left: float, top: float, width: float, height: float, count: int
) -> list[tuple[float, float, float, float]]:
    """Layout boxes for full-image slides (supports up to 4 images)."""
    gutter = 0.25
    if count <= 1:
        return [(left, top, width, height)]
    if count == 2:
        w = (width - gutter) / 2
        return [
            (left, top, w, height),
            (left + w + gutter, top, w, height),
        ]
    if count == 3:
        w = (width - gutter) / 2
        h = (height - gutter) / 2
        return [
            (left, top, w, height),
            (left + w + gutter, top, w, h),
            (left + w + gutter, top + h + gutter, w, h),
        ]

    w = (width - gutter) / 2
    h = (height - gutter) / 2
    return [
        (left, top, w, h),
        (left + w + gutter, top, w, h),
        (left, top + h + gutter, w, h),
        (left + w + gutter, top + h + gutter, w, h),
    ]


def add_side_images(
    slide,
    image_paths: list[Path],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    max_images: int = 3,
) -> None:
    """Add up to `max_images` stacked images in the given box (in inches)."""
    if width <= 0 or height <= 0:
        return

    paths = image_paths[:max_images]
    if not paths:
        return

    gutter = 0.2
    per_height = (height - gutter * (len(paths) - 1)) / len(paths)
    current_top = top
    for image_path in paths:
        add_picture_fitted(slide, image_path, left, current_top, width, per_height)
        current_top += per_height + gutter


def set_slide_notes(slide, notes: str) -> None:
    """Set speaker notes for a slide (best-effort)."""
    note_lines = prepare_lines(notes or "")
    text = "\n".join(line for line in note_lines if line.strip())
    if not text:
        return
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        return


def add_picture_fitted(
    slide, image_path: Path, left: float, top: float, width: float, height: float
) -> None:
    """Add an image to the slide, preserving aspect ratio within the given box (in inches)."""
    if width <= 0 or height <= 0:
        return

    dims = get_image_dimensions_px(image_path)
    if not dims:
        slide.shapes.add_picture(
            str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height)
        )
        return

    img_w_px, img_h_px = dims
    if img_w_px <= 0 or img_h_px <= 0:
        slide.shapes.add_picture(
            str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height)
        )
        return

    img_aspect = img_w_px / img_h_px
    box_aspect = width / height

    if img_aspect >= box_aspect:
        fitted_width = width
        fitted_height = width / img_aspect
        fitted_left = left
        fitted_top = top + (height - fitted_height) / 2
    else:
        fitted_height = height
        fitted_width = height * img_aspect
        fitted_left = left + (width - fitted_width) / 2
        fitted_top = top

    slide.shapes.add_picture(
        str(image_path),
        Inches(fitted_left),
        Inches(fitted_top),
        width=Inches(fitted_width),
        height=Inches(fitted_height),
    )


def get_image_dimensions_px(path: Path) -> tuple[int, int] | None:
    """Best-effort image dimension extraction (PNG/JPEG/GIF)."""
    try:
        data = path.read_bytes()[:512_000]
    except OSError:
        return None

    if len(data) < 10:
        return None

    png = _get_png_dimensions(data)
    if png:
        return png
    gif = _get_gif_dimensions(data)
    if gif:
        return gif
    jpeg = _get_jpeg_dimensions(data)
    if jpeg:
        return jpeg
    return None


def _get_png_dimensions(data: bytes) -> tuple[int, int] | None:
    if len(data) < 24:
        return None
    if not data.startswith(b"\x89PNG\r\n\x1a\n"):
        return None
    if data[12:16] != b"IHDR":
        return None
    try:
        width, height = struct.unpack(">II", data[16:24])
    except struct.error:
        return None
    return int(width), int(height)


def _get_gif_dimensions(data: bytes) -> tuple[int, int] | None:
    if len(data) < 10:
        return None
    if data[:6] not in {b"GIF87a", b"GIF89a"}:
        return None
    try:
        width, height = struct.unpack("<HH", data[6:10])
    except struct.error:
        return None
    return int(width), int(height)


def _get_jpeg_dimensions(data: bytes) -> tuple[int, int] | None:
    if len(data) < 4 or data[0:2] != b"\xff\xd8":
        return None

    i = 2
    length = len(data)
    while i < length:
        if data[i] != 0xFF:
            i += 1
            continue
        while i < length and data[i] == 0xFF:
            i += 1
        if i >= length:
            break
        marker = data[i]
        i += 1

        if marker in {0xD8, 0xD9}:
            continue
        if marker == 0xDA:  # SOS
            break
        if i + 2 > length:
            break

        try:
            seg_len = struct.unpack(">H", data[i : i + 2])[0]
        except struct.error:
            return None
        if seg_len < 2:
            return None

        sof_markers = {
            0xC0,
            0xC1,
            0xC2,
            0xC3,
            0xC5,
            0xC6,
            0xC7,
            0xC9,
            0xCA,
            0xCB,
            0xCD,
            0xCE,
            0xCF,
        }
        if marker in sof_markers and i + 7 <= length:
            try:
                height = struct.unpack(">H", data[i + 3 : i + 5])[0]
                width = struct.unpack(">H", data[i + 5 : i + 7])[0]
            except struct.error:
                return None
            return int(width), int(height)

        i += seg_len

    return None
