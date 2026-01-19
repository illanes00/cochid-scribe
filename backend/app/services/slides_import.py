"""Import PPTX slides into slides_data with asset extraction."""

from __future__ import annotations

import html
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from sqlalchemy.orm import Session

from app.models.asset import Asset

DEFAULT_THEME = {
    "primaryColor": "#1a365d",
    "secondaryColor": "#c53030",
    "fontFamily": "IBM Plex Sans, system-ui, sans-serif",
}


def _escape(text: str) -> str:
    return html.escape(text, quote=True)


def _store_image(
    image_blob: bytes,
    ext: str,
    document_id: str,
    upload_dir: Path,
    db: Session,
) -> str:
    asset_id = str(uuid.uuid4())
    safe_ext = f".{ext.lower().lstrip('.')}" if ext else ""
    filename = f"{asset_id}{safe_ext}"
    filepath = upload_dir / filename
    filepath.write_bytes(image_blob)

    asset = Asset(
        id=asset_id,
        document_id=document_id,
        filename=filename,
        mime_type=f"image/{ext.lower()}" if ext else "image/png",
        size_bytes=len(image_blob),
        url=f"/uploads/assets/{filename}",
    )
    db.add(asset)
    return asset.url


def build_slides_data_from_pptx(
    pptx_path: Path,
    db: Session,
    document_id: str,
    upload_dir: Path,
) -> dict:
    """Parse a PPTX file into slides_data and persist images as assets."""
    upload_dir.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(str(pptx_path))
    slides = []

    for index, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title else f"Slide {index}"
        body_lines: list[str] = []
        image_urls: list[str] = []

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                url = _store_image(image.blob, image.ext, document_id, upload_dir, db)
                image_urls.append(url)
                continue
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        body_lines.append(text)

        content_parts = [f"<p>{_escape(line)}</p>" for line in body_lines]
        for url in image_urls:
            content_parts.append(f'<img src="{url}" alt="Slide image" />')

        layout = "image-full" if image_urls and not body_lines else "content"

        slides.append(
            {
                "id": f"slide-{index}",
                "slideNumber": index,
                "layout": layout,
                "title": title,
                "content": "".join(content_parts),
                "notes": "",
            }
        )

    return {"slides": slides, "theme": DEFAULT_THEME}
