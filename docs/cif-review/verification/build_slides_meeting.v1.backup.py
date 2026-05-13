"""Build PPT para reunión Carla/Eduardo del jueves 7 may 2026.

12 slides con identidad visual EP (paleta health serie isapres-2023).
Genera output/slides/meeting-2026-05-07.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

BASE = Path("/srv/projects/cochid/cochid-scribe/docs/cif-review")
OUT_DIR = BASE / "output/slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT = OUT_DIR / "meeting-2026-05-07.pptx"

# Paleta EP serie health
COL_PRIMARY = RGBColor(0x60, 0x40, 0x80)   # purple
COL_SECONDARY = RGBColor(0x80, 0x40, 0x60) # plum
COL_ACCENT = RGBColor(0xC0, 0x60, 0x20)    # warm orange
COL_DARK = RGBColor(0x1A, 0x36, 0x5D)      # navy
COL_TEXT = RGBColor(0x1A, 0x20, 0x2C)      # near black
COL_MUTED = RGBColor(0x4A, 0x55, 0x68)     # slate
COL_BG_SOFT = RGBColor(0xFA, 0xF5, 0xF0)   # warm off-white
COL_RULE = RGBColor(0xCB, 0xD5, 0xE0)      # light gray
COL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEADING = "Calibri"   # weasyprint usa Barlow/Nunito; ppt usa lo que tenga el sistema
FONT_BODY = "Calibri"


def add_text(slide, left, top, width, height, text, size, bold=False,
             color=COL_TEXT, font=FONT_BODY, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, size=18, color=COL_TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•   " + b
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(8)


def add_rect(slide, left, top, width, height, color):
    """Add a solid rectangle (shape type 1)."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_header_band(slide, color=COL_PRIMARY, height=Inches(0.85)):
    """Top color band for content slides."""
    add_rect(slide, 0, 0, prs.slide_width, height, color)


def add_footer_brand(slide):
    """Footer with EP brand on the right."""
    add_text(slide, Inches(11.5), Inches(7.0), Inches(1.7), Inches(0.4),
             "ESPACIO PÚBLICO", 8, bold=True, color=COL_MUTED,
             font=FONT_HEADING, align=PP_ALIGN.RIGHT)


def slide_title(slide, title_text):
    """Standard slide title in header band."""
    add_header_band(slide)
    add_text(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.5),
             title_text, 26, bold=True, color=COL_WHITE, font=FONT_HEADING)
    add_text(slide, Inches(0.6), Inches(0.62), Inches(12), Inches(0.25),
             "INCLUSIÓN SOSTENIBLE DE MEDICAMENTOS — SESIÓN DE CIERRE EDITORIAL · 7 MAYO 2026",
             8, color=RGBColor(0xE0, 0xE0, 0xF0), font=FONT_HEADING)


# ============================================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ----- Slide 1: Title -----
s = prs.slides.add_slide(blank)
# Gradient simulado: 4 bandas verticales orange->plum->purple
band_h = prs.slide_height // 4
add_rect(s, 0, 0, prs.slide_width, band_h, RGBColor(0xE0, 0x80, 0x00))
add_rect(s, 0, band_h, prs.slide_width, band_h, RGBColor(0xC0, 0x60, 0x20))
add_rect(s, 0, 2*band_h, prs.slide_width, band_h, RGBColor(0x80, 0x40, 0x60))
add_rect(s, 0, 3*band_h, prs.slide_width, band_h, RGBColor(0x60, 0x40, 0x80))

add_text(s, Inches(1), Inches(1.0), Inches(11.333), Inches(0.5),
         "MINUTA TÉCNICA · CIERRE EDITORIAL", 14, bold=True,
         color=COL_WHITE, font=FONT_HEADING, align=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(2.4), Inches(11.333), Inches(2.2),
         "Inclusión sostenible de medicamentos\nen los planes de salud en Chile",
         44, bold=True, color=COL_WHITE, font=FONT_HEADING, align=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(5.0), Inches(11.333), Inches(0.4),
         "Sesión de cierre editorial con directores", 18,
         color=COL_WHITE, font=FONT_BODY, align=PP_ALIGN.CENTER, italic=True)

add_text(s, Inches(1), Inches(5.6), Inches(11.333), Inches(0.4),
         "Carla Castillo · Eduardo Undurraga · Eleni Kokkidou · Benjamín García · Martín Illanes",
         13, color=COL_WHITE, font=FONT_BODY, align=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(6.7), Inches(11.333), Inches(0.4),
         "ESPACIO PÚBLICO · 7 DE MAYO DE 2026", 11, bold=True,
         color=COL_WHITE, font=FONT_HEADING, align=PP_ALIGN.CENTER)


# ----- Slide 2: Estado del informe -----
s = prs.slides.add_slide(blank)
slide_title(s, "Estado del informe")
add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "Versión 3.17 (mayo 5)", 22, bold=True, color=COL_PRIMARY)
add_bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(2.5), [
    "57 páginas · 9 tablas · 33 figuras · 8 capítulos + 6 anexos",
    "v3.16 enviada el 24 de abril; v3.17 con cierres cosméticos posteriores",
    "Tabla de protección completa: GES + LRS + DAC + MLE + Arsenal APS + Receta cautiva",
    "Foco comparado en OCDE europeo + Uruguay (FNR), justificación metodológica",
    "Innovación con valor sanitario incorporada como agenda explícita",
    "Tres escenarios de política con trade-offs cuantificados",
])
add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
         "Calendario hacia publicación", 16, bold=True, color=COL_PRIMARY)
add_bullets(s, Inches(0.6), Inches(5.9), Inches(12), Inches(1.2), [
    "Hoy 5 may: brief 8pp + PPT + v3.17 listos para esta sesión",
    "8–15 may: incorporación de feedback de directores → v3.18",
    "20–25 may: cierre y publicación",
])
add_footer_brand(s)


# ----- Slide 3: El panorama -----
s = prs.slides.add_slide(blank)
slide_title(s, "El panorama")
add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "Aseguramiento amplio en el papel, protección financiera estrecha en la práctica",
         18, bold=True, color=COL_SECONDARY, italic=True)
add_bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(3.5), [
    "Hogares financian directamente cerca del 71 % del gasto retail farmacéutico (OECD SHA, 2022)",
    "Promedio OCDE: 39 %",
    "Tratamientos de alto costo siguen accediéndose vía judicialización cuando las garantías explícitas no los cubren",
    "El sistema funciona, pero no como debiera",
])
add_rect(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.7), COL_BG_SOFT)
add_text(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.4),
         "El brief ordena la conversación", 16, bold=True, color=COL_PRIMARY)
add_bullets(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(1.1), [
    "Diagnóstico de las dos lógicas del problema (acumulativa + catastrófica)",
    "Mapa completo de la tabla de protección actual",
    "Evidencia comparada con foco europeo",
    "Tres escenarios con trade-offs explícitos",
], size=14)
add_footer_brand(s)


# ----- Slide 4: Mensajes clave -----
s = prs.slides.add_slide(blank)
slide_title(s, "Mensajes clave")

# Cinco bullets en dos columnas
mensajes = [
    ("1", "El gasto de bolsillo es estructural, no circunstancial",
     "71 % del gasto retail farmacéutico financiado por hogares (cerca del doble del promedio OCDE)"),
    ("2", "El problema tiene dos lógicas distintas y simultáneas",
     "Riesgo acumulativo (crónicos ambulatorios) + riesgo catastrófico (alto costo). Cada lógica requiere instrumentos diferentes."),
    ("3", "La tabla de protección actual es más amplia que GES + LRS",
     "Incluye DAC, MLE, Arsenal APS y receta cautiva. La fragmentación, no la ausencia, explica buena parte de la brecha."),
    ("4", "La judicialización es válvula institucional, no fenómeno lateral",
     "Síntoma de falla regulatoria; reaparece incluso en países con cobertura universal nominal."),
    ("5", "La evidencia comparada apunta a un patrón de paquete",
     "Canasta explícita + copagos con topes + dispensación con convenio. Ningún componente aislado ha resuelto el problema."),
]

y = Inches(1.2)
for num, title, body in mensajes:
    # Numero grande naranja
    add_text(s, Inches(0.6), y, Inches(0.6), Inches(1.0),
             num, 36, bold=True, color=COL_ACCENT, font=FONT_HEADING)
    add_text(s, Inches(1.3), y + Emu(50000), Inches(11.5), Inches(0.4),
             title, 15, bold=True, color=COL_PRIMARY)
    add_text(s, Inches(1.3), y + Inches(0.45), Inches(11.5), Inches(0.5),
             body, 12, color=COL_TEXT)
    y += Inches(1.15)

add_footer_brand(s)


# ----- Slide 5: Diagnóstico dual -----
s = prs.slides.add_slide(blank)
slide_title(s, "Diagnóstico: dos lógicas, dos instrumentos")

# Dos columnas
add_rect(s, Inches(0.6), Inches(1.3), Inches(6), Inches(5.5), COL_BG_SOFT)
add_rect(s, Inches(6.8), Inches(1.3), Inches(6), Inches(5.5), COL_BG_SOFT)

add_text(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.6),
         "(1) Lógica acumulativa", 22, bold=True, color=COL_ACCENT, font=FONT_HEADING)
add_bullets(s, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.2), [
    "Crónicos ambulatorios: hipertensión, diabetes, salud mental, asma",
    "Gasto recurrente que se acumula mes a mes",
    "29 % dejó de tomar dosis por costo (Ipsos-EP, jul 2025)",
    "Ausencia de tope anual = brecha estructural visible",
    "Subsidio uniforme al precio (IVA) atiende parcialmente, no resuelve",
], size=14)

add_text(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.6),
         "(2) Lógica catastrófica", 22, bold=True, color=COL_PRIMARY, font=FONT_HEADING)
add_bullets(s, Inches(7.0), Inches(2.2), Inches(5.6), Inches(4.2), [
    "Oncológicos modernos, enfermedades poco frecuentes, biotecnológicos",
    "Riesgo de agotar patrimonio familiar",
    "Sistema chileno: GES + LRS + DAC + judicialización (fragmentado)",
    "Ley Ricarte Soto: lista cerrada, ingresos discrecionales",
    "Cobertura universal nominal no elimina judicialización (cf. Colombia)",
], size=14)

add_footer_brand(s)


# ----- Slide 6: Tabla de protección completa -----
s = prs.slides.add_slide(blank)
slide_title(s, "Tabla de protección farmacéutica vigente")

instrumentos = [
    ("GES", "FONASA + ISAPRE", "87 patologías, incluye ambulatorios para parte de ellas"),
    ("Ley Ricarte Soto", "FONASA + ISAPRE", "Diagnósticos y tratamientos de alto costo, lista cerrada"),
    ("DAC", "FONASA", "Drogas oncológicas y de alto costo no GES/LRS"),
    ("Modalidad Libre Elección", "FONASA tramos B–D", "Bonificación parcial; no crónicos generales"),
    ("Arsenal Farmacológico APS", "FONASA en APS", "Esenciales con dispensación gratuita; stock heterogéneo"),
    ("Receta cautiva", "FONASA", "Continuidad de tratamientos hospitalarios; ligada al prestador"),
]

# Header
y = Inches(1.3)
add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.45), COL_PRIMARY)
add_text(s, Inches(0.7), y + Emu(40000), Inches(3.5), Inches(0.35),
         "INSTRUMENTO", 12, bold=True, color=COL_WHITE, font=FONT_HEADING)
add_text(s, Inches(4.3), y + Emu(40000), Inches(3.0), Inches(0.35),
         "BENEFICIARIOS", 12, bold=True, color=COL_WHITE, font=FONT_HEADING)
add_text(s, Inches(7.4), y + Emu(40000), Inches(5.3), Inches(0.35),
         "TIPO DE COBERTURA", 12, bold=True, color=COL_WHITE, font=FONT_HEADING)
y += Inches(0.55)

for i, (inst, ben, cob) in enumerate(instrumentos):
    if i % 2 == 0:
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.7), COL_BG_SOFT)
    add_text(s, Inches(0.7), y + Inches(0.1), Inches(3.5), Inches(0.6),
             inst, 12, bold=True, color=COL_SECONDARY)
    add_text(s, Inches(4.3), y + Inches(0.1), Inches(3.0), Inches(0.6),
             ben, 11, color=COL_TEXT)
    add_text(s, Inches(7.4), y + Inches(0.1), Inches(5.3), Inches(0.6),
             cob, 11, color=COL_TEXT)
    y += Inches(0.78)

add_footer_brand(s)


# ----- Slide 7: Evidencia comparada -----
s = prs.slides.add_slide(blank)
slide_title(s, "Lo que muestra la evidencia comparada")
add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "Foco OCDE europeo: Australia, España, Francia, Inglaterra, Alemania, Portugal + Uruguay (FNR)",
         15, italic=True, color=COL_MUTED)

# Tabla comparada
y = Inches(1.85)
headers = ["País", "Canasta", "Topes anuales", "Dispensación", "Pago por valor"]
widths_in = [3.0, 2.0, 2.0, 2.5, 2.5]
xs = [Inches(0.6)]
for w in widths_in[:-1]:
    xs.append(xs[-1] + Inches(w))

add_rect(s, Inches(0.6), y, Inches(12), Inches(0.45), COL_PRIMARY)
for i, (h, w) in enumerate(zip(headers, widths_in)):
    add_text(s, xs[i] + Inches(0.05), y + Emu(40000), Inches(w), Inches(0.35),
             h, 12, bold=True, color=COL_WHITE, font=FONT_HEADING)
y += Inches(0.55)

countries = [
    ("Australia (PBS)", "✓", "✓", "✓", "✓"),
    ("España", "✓", "✓", "✓", "parcial"),
    ("Francia", "✓", "✓", "✓", "parcial"),
    ("Inglaterra (NHS)", "✓", "✓", "✓", "✓"),
    ("Alemania", "✓", "✓", "✓", "parcial"),
    ("Uruguay (FNR)", "✓", "✓", "✓", "parcial"),
    ("Chile (actual)", "GES+LRS+DAC", "—", "limitada", "limitada"),
]
for i, row in enumerate(countries):
    bg = COL_BG_SOFT if i % 2 == 0 else COL_WHITE
    add_rect(s, Inches(0.6), y, Inches(12), Inches(0.5), bg)
    is_chile = row[0].startswith("Chile")
    color = COL_ACCENT if is_chile else COL_TEXT
    bold = is_chile
    for j, (val, w) in enumerate(zip(row, widths_in)):
        add_text(s, xs[j] + Inches(0.05), y + Inches(0.1), Inches(w), Inches(0.4),
                 val, 12, bold=bold, color=color)
    y += Inches(0.55)

add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
         "Patrón consistente: tres componentes simultáneos. Ningún componente aislado ha resuelto el problema.",
         13, italic=True, bold=True, color=COL_PRIMARY)

add_footer_brand(s)


# ----- Slide 8: Tres escenarios -----
s = prs.slides.add_slide(blank)
slide_title(s, "Tres escenarios de política para Chile")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "Trade-offs explícitos. El informe no zanja preferencias; ordena la conversación.",
         14, italic=True, color=COL_MUTED)

escenarios = [
    {"name": "E1 · Cobertura focalizada ampliada", "color": RGBColor(0x80, 0xA0, 0xC0),
     "rows": [("Lógica", "Profundizar instrumentos vigentes"),
              ("Reducción OOP", "71 % → 60 %"),
              ("Costo fiscal", "USD 200–400 M/año"),
              ("Reforma legal", "Modificación GES, LRS, DAC")]},
    {"name": "E2 · Beneficio Farmacéutico Universal", "color": RGBColor(0x80, 0x40, 0x60),
     "rows": [("Lógica", "Reordenar cobertura ambulatoria, reglas comunes"),
              ("Reducción OOP", "71 % → 45–50 %"),
              ("Costo fiscal", "USD 800–1.200 M/año"),
              ("Reforma legal", "Ley marco BFAU")]},
    {"name": "E3 · Convergencia plena OCDE", "color": RGBColor(0x60, 0x40, 0x80),
     "rows": [("Lógica", "Tope anual universal de gasto OOP"),
              ("Reducción OOP", "71 % → 25–30 %"),
              ("Costo fiscal", "USD 2.500 M+/año"),
              ("Reforma legal", "Reforma sistémica integrada con FONASA")]},
]

x = Inches(0.6)
col_w = Inches(4.05)
y_top = Inches(1.85)
for esc in escenarios:
    add_rect(s, x, y_top, col_w, Inches(0.5), esc["color"])
    add_text(s, x + Inches(0.15), y_top + Inches(0.1), col_w, Inches(0.4),
             esc["name"], 13, bold=True, color=COL_WHITE, font=FONT_HEADING)
    yy = y_top + Inches(0.55)
    for k, v in esc["rows"]:
        add_text(s, x + Inches(0.15), yy, col_w - Inches(0.3), Inches(0.3),
                 k.upper(), 9, bold=True, color=COL_MUTED, font=FONT_HEADING)
        add_text(s, x + Inches(0.15), yy + Inches(0.28), col_w - Inches(0.3), Inches(0.7),
                 v, 12, color=COL_TEXT)
        yy += Inches(0.95)
    x += col_w + Inches(0.15)

add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
         "El BFAU (E2) se desarrolla con mayor detalle por riqueza analítica, no por preferencia.",
         12, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER)

add_footer_brand(s)


# ----- Slide 9: Innovación con valor -----
s = prs.slides.add_slide(blank)
slide_title(s, "Innovación con valor sanitario")

add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "Agenda explícita transversal a los tres escenarios", 16,
         italic=True, color=COL_MUTED)

innov = [
    ("Empagliflozina", "Diabetes T2 con riesgo CV", "Reducción mortalidad CV", "Zinman, NEJM 2015"),
    ("Dapagliflozina", "Insuficiencia cardíaca FE reducida", "Reducción muerte CV y hospitalización", "McMurray, NEJM 2019"),
    ("Trikafta (CFTR)", "Fibrosis quística", "Cambio en historia natural", "Middleton, NEJM 2019"),
    ("Pembrolizumab", "Melanoma metastásico", "Sobrevida superior, respuestas duraderas", "Schadendorf, NEJM 2015"),
    ("Onasemnogén", "AME tipo 1", "Una dosis, cambio sustancial", "Mendell, Nat Med 2022"),
]

y = Inches(2.0)
add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.45), COL_PRIMARY)
add_text(s, Inches(0.7), y + Emu(40000), Inches(2.6), Inches(0.35),
         "INNOVACIÓN", 11, bold=True, color=COL_WHITE, font=FONT_HEADING)
add_text(s, Inches(3.4), y + Emu(40000), Inches(2.8), Inches(0.35),
         "INDICACIÓN", 11, bold=True, color=COL_WHITE, font=FONT_HEADING)
add_text(s, Inches(6.4), y + Emu(40000), Inches(3.6), Inches(0.35),
         "APORTE SANITARIO", 11, bold=True, color=COL_WHITE, font=FONT_HEADING)
add_text(s, Inches(10.2), y + Emu(40000), Inches(2.5), Inches(0.35),
         "REFERENTE", 11, bold=True, color=COL_WHITE, font=FONT_HEADING)
y += Inches(0.55)

for i, (inn, ind, ap, ref) in enumerate(innov):
    if i % 2 == 0:
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.65), COL_BG_SOFT)
    add_text(s, Inches(0.7), y + Inches(0.1), Inches(2.6), Inches(0.55),
             inn, 11, bold=True, color=COL_SECONDARY)
    add_text(s, Inches(3.4), y + Inches(0.1), Inches(2.8), Inches(0.55),
             ind, 10, color=COL_TEXT)
    add_text(s, Inches(6.4), y + Inches(0.1), Inches(3.6), Inches(0.55),
             ap, 10, color=COL_TEXT)
    add_text(s, Inches(10.2), y + Inches(0.1), Inches(2.5), Inches(0.55),
             ref, 10, italic=True, color=COL_MUTED)
    y += Inches(0.7)

add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
         "Mecanismos: ETESA orientada a valor + acuerdos de riesgo compartido para terapias de alto costo.",
         13, italic=True, color=COL_PRIMARY, align=PP_ALIGN.CENTER)

add_footer_brand(s)


# ----- Slide 10: Triple-pack -----
s = prs.slides.add_slide(blank)
slide_title(s, "Lo que entrego en esta sesión")

add_rect(s, Inches(0.6), Inches(1.5), Inches(4.0), Inches(5.5), COL_BG_SOFT)
add_rect(s, Inches(4.7), Inches(1.5), Inches(4.0), Inches(5.5), COL_BG_SOFT)
add_rect(s, Inches(8.8), Inches(1.5), Inches(4.0), Inches(5.5), COL_BG_SOFT)

# Long doc
add_text(s, Inches(0.8), Inches(1.7), Inches(3.6), Inches(0.5),
         "INFORME EXTENSO", 14, bold=True, color=COL_PRIMARY, font=FONT_HEADING)
add_text(s, Inches(0.8), Inches(2.1), Inches(3.6), Inches(0.5),
         "v3.17", 32, bold=True, color=COL_SECONDARY, font=FONT_HEADING)
add_bullets(s, Inches(0.8), Inches(2.9), Inches(3.6), Inches(4.0), [
    "57 páginas",
    "8 capítulos + 6 anexos",
    "9 tablas, 33 figuras",
    "Cierre cosmético v3.16 → v3.17",
    "Tabla protección completa",
    "Foco UE justificado",
    "Innovación con valor",
    "Tracked changes preservados",
], size=12)

# Brief
add_text(s, Inches(4.9), Inches(1.7), Inches(3.6), Inches(0.5),
         "POLICY BRIEF", 14, bold=True, color=COL_PRIMARY, font=FONT_HEADING)
add_text(s, Inches(4.9), Inches(2.1), Inches(3.6), Inches(0.5),
         "8 pp", 32, bold=True, color=COL_SECONDARY, font=FONT_HEADING)
add_bullets(s, Inches(4.9), Inches(2.9), Inches(3.6), Inches(4.0), [
    "Síntesis para política pública",
    "Identidad visual EP",
    "Panorama → opciones",
    "5 mensajes clave",
    "Tres escenarios",
    "Innovación con valor",
    "Nota institucional concisa",
    "Para CIF + público",
], size=12)

# PPT
add_text(s, Inches(9.0), Inches(1.7), Inches(3.6), Inches(0.5),
         "ESTA PRESENTACIÓN", 14, bold=True, color=COL_PRIMARY, font=FONT_HEADING)
add_text(s, Inches(9.0), Inches(2.1), Inches(3.6), Inches(0.5),
         "12 sl.", 32, bold=True, color=COL_SECONDARY, font=FONT_HEADING)
add_bullets(s, Inches(9.0), Inches(2.9), Inches(3.6), Inches(4.0), [
    "Para esta sesión deliberativa",
    "Ancla la discusión",
    "Resume el panorama",
    "Presenta opciones",
    "Espacio para feedback",
    "No es la presentación pública",
], size=12)

add_footer_brand(s)


# ----- Slide 11: Próximas decisiones (placeholder) -----
s = prs.slides.add_slide(blank)
slide_title(s, "Próximas decisiones")

add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
         "Puntos sobre los que necesito su mirada", 24,
         bold=True, color=COL_PRIMARY, font=FONT_HEADING)

add_bullets(s, Inches(0.6), Inches(2.1), Inches(12), Inches(4.5), [
    "Cuán prescriptivo debe ser el informe — tensión entre profundidad técnica y preservar rol de honest broker",
    "BFAU bajo arquitectura FONASA-ISAPRE actual, o subordinado a reforma sistémica pendiente",
    "Tratamiento de la judicialización — promesa de reducción o señal estructural permanente",
    "Foco geográfico — OCDE europeo puro o incorporación de casos LatAm con sus problemas",
    "Institucionalidad de evaluación e incorporación — agencia técnica nueva o estructuras existentes; RSA opt-in u obligatorios",
], size=15)

add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "Estas son decisiones de marco conceptual. Las decisiones operativas las puedo calibrar.",
         13, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER)

add_footer_brand(s)


# ----- Slide 12: Próximos pasos -----
s = prs.slides.add_slide(blank)
slide_title(s, "Próximos pasos hacia la publicación")

# Timeline horizontal
y = Inches(2.0)
hitos = [
    ("HOY", "Sesión", "Brief + PPT + v3.17"),
    ("8–15 may", "v3.18", "Incorporación feedback directores"),
    ("15–22 may", "Revisión final", "Eleni + Benja"),
    ("25 may", "Publicación", "Brief PDF + Informe extenso"),
    ("Q3 2026", "Seminario", "Discusión pública"),
]
x = Inches(0.6)
hito_w = Inches(2.4)
for i, (date, milestone, detail) in enumerate(hitos):
    color = COL_PRIMARY if i < len(hitos) - 1 else COL_ACCENT
    add_rect(s, x, y, hito_w, Inches(0.6), color)
    add_text(s, x + Inches(0.1), y + Inches(0.1), hito_w - Inches(0.2), Inches(0.4),
             date, 14, bold=True, color=COL_WHITE, font=FONT_HEADING, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(0.7), hito_w - Inches(0.2), Inches(0.4),
             milestone, 14, bold=True, color=COL_PRIMARY, font=FONT_HEADING, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(1.15), hito_w - Inches(0.2), Inches(1.0),
             detail, 11, color=COL_TEXT, align=PP_ALIGN.CENTER)
    x += hito_w + Inches(0.05)

add_text(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.5),
         "Gracias", 36, bold=True, color=COL_PRIMARY, font=FONT_HEADING, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
         "Quedo atento al feedback que surja en esta sesión y en los próximos días.",
         14, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER)

add_footer_brand(s)


# ----- Save -----
prs.save(str(OUT))
print(f"PPTX saved: {OUT}")
print(f"Size: {OUT.stat().st_size:,} bytes")
print(f"Slides: {len(prs.slides)}")
