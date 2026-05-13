"""Build PPT para reunión Carla/Eduardo del jueves 7 may 2026.

12 slides con identidad visual EP (paleta health serie isapres-2023).
Genera output/slides/meeting-2026-05-07.pptx

v4: alineado con informe v4.3. Slide 9 reemplazado: zoom BFU.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

BASE = Path("/srv/projects/cochid/cochid-scribe/docs/cif-review")
OUT_DIR = BASE / "output/slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT = OUT_DIR / "meeting-2026-05-07.pptx"

# Paleta EP serie health
COL_PRIMARY = RGBColor(0x60, 0x40, 0x80)   # purple
COL_SECONDARY = RGBColor(0x80, 0x40, 0x60) # plum
COL_ACCENT = RGBColor(0xC0, 0x60, 0x20)    # warm orange
COL_DARK = RGBColor(0x1A, 0x36, 0x5D)      # navy
COL_TEXT = RGBColor(0x1A, 0x20, 0x2C)      # near black
COL_MUTED = RGBColor(0x4A, 0x55, 0x68)     # slate
COL_BG_SOFT = RGBColor(0xFA, 0xF5, 0xF0)   # warm off-white
COL_RULE = RGBColor(0xCB, 0xD5, 0xE0)      # light gray
COL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEADING = "Calibri"   # weasyprint usa Barlow/Nunito; ppt usa lo que tenga el sistema
FONT_BODY = "Calibri"


def add_text(slide, left, top, width, height, text, size, bold=False,
             color=COL_TEXT, font=FONT_BODY, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, size=18, color=COL_TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•   " + b
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(8)


def add_rect(slide, left, top, width, height, color):
    """Add a solid rectangle (shape type 1)."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_header_band(slide, color=COL_PRIMARY, height=Inches(0.85)):
    """Top color band for content slides."""
    add_rect(slide, 0, 0, prs.slide_width, height, color)


def add_footer_brand(slide):
    """Footer with EP brand on the right."""
    add_text(slide, Inches(11.5), Inches(7.0), Inches(1.7), Inches(0.4),
             "ESPACIO PÚBLICO", 8, bold=True, color=COL_MUTED,
             font=FONT_HEADING, align=PP_ALIGN.RIGHT)


def slide_title(slide, title_text):
    """Standard slide title in header band."""
    add_header_band(slide)
    add_text(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.5),
             title_text, 26, bold=True, color=COL_WHITE, font=FONT_HEADING)
    add_text(slide, Inches(0.6), Inches(0.62), Inches(12), Inches(0.25),
             "INCLUSIÓN SOSTENIBLE DE MEDICAMENTOS · SESIÓN DE CIERRE EDITORIAL · 7 MAYO 2026",
             8, color=RGBColor(0xE0, 0xE0, 0xF0), font=FONT_HEADING)


# ============================================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ----- Slide 1: Title -----
s = prs.slides.add_slide(blank)
# Gradient simulado: 4 bandas verticales orange->plum->purple
band_h = prs.slide_height // 4
add_rect(s, 0, 0, prs.slide_width, band_h, RGBColor(0xE0, 0x80, 0x00))
add_rect(s, 0, band_h, prs.slide_width, band_h, RGBColor(0xC0, 0x60, 0x20))
add_rect(s, 0, 2*band_h, prs.slide_width, band_h, RGBColor(0x80, 0x40, 0x60))
add_rect(s, 0, 3*band_h, prs.slide_width, band_h, RGBColor(0x60, 0x40, 0x80))

add_text(s, Inches(1), Inches(1.0), Inches(11.333), Inches(0.5),
         "MINUTA TÉCNICA · CIERRE EDITORIAL", 14, bold=True,
         color=COL_WHITE, font=FONT_HEADING, align=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(2.4), Inches(11.333), Inches(2.2),
         "Inclusión sostenible de medicamentos\nen los planes de salud en Chile",
         44, bold=True, color=COL_WHITE, font=FONT_HEADING, align=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(5.0), Inches(11.333), Inches(0.4),
         "Sesión de cierre editorial con directores", 18,
         color=COL_WHITE, font=FONT_BODY, align=PP_ALIGN.CENTER, italic=True)

add_text(s, Inches(1), Inches(5.6), Inches(11.333), Inches(0.4),
         "Carla Castillo · Eduardo Undurraga · Eleni Kokkidou · Benjamín García · Martín Illanes",
         13, color=COL_WHITE, font=FONT_BODY, align=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(6.7), Inches(11.333), Inches(0.4),
         "ESPACIO PÚBLICO · 7 DE MAYO DE 2026", 11, bold=True,
         color=COL_WHITE, font=FONT_HEADING, align=PP_ALIGN.CENTER)


# ----- Slide 2: Estado del informe -----
s = prs.slides.add_slide(blank)
slide_title(s, "Estado del informe")
add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "Versión v4.3 (mayo 7)", 22, bold=True, color=COL_PRIMARY)
add_bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(2.5), [
    "57 páginas · 9 tablas · 33 figuras · 8 capítulos + 6 anexos",
    "Marco BFU como zoom analítico para la discusión",
    "Cifras dobles: 62 % gasto total / 71 % retail (CIF/UC 2024 y OECD SHA 2022)",
    "Gasto público: 0,46 % PIB total · 0,37 % PIB en HC51 (2023)",
    "Tabla de protección con 6 instrumentos (GES, LRS, DAC, FOFAR+APS, CAEC, Cenabast)",
    "Tres escenarios cuantificados con BFU (E2) como zoom analítico",
], size=15)
add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
         "Calendario hacia publicación", 16, bold=True, color=COL_PRIMARY)
add_bullets(s, Inches(0.6), Inches(5.9), Inches(12), Inches(1.2), [
    "Hoy 7 may: brief v4.3 + PPT v4 + informe v4.3 listos para esta sesión",
    "8 a 20 may: incorporación de feedback de Carla Castillo y Eduardo Undurraga",
    "20-25 may: revisión Eleni + Benja",
    "25 may: publicación",
], size=14)
add_footer_brand(s)


# ----- Slide 3: El panorama -----
s = prs.slides.add_slide(blank)
slide_title(s, "El panorama")
add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "Aseguramiento amplio en el papel, protección financiera estrecha en la práctica",
         18, bold=True, color=COL_SECONDARY, italic=True)
add_bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(3.5), [
    "62 % del gasto total en medicamentos (CIF/UC, 2024) y 71 % del retail (OECD SHA, 2022) lo financian los hogares",
    "Gasto público: 0,46 % PIB total · 0,37 % PIB en HC51 (2023)",
    "21,8 % de los hogares gastan más del 10 % de su ingreso per cápita en medicamentos",
    "Tratamientos de alto costo accedidos vía judicialización ($81.000 M en 2024)",
])
add_rect(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.7), COL_BG_SOFT)
add_text(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.4),
         "El brief ordena la conversación", 16, bold=True, color=COL_PRIMARY)
add_bullets(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(1.1), [
    "Diagnóstico dual: mirada acumulativa y mirada catastrófica",
    "Mapa completo de la tabla de protección actual",
    "Evidencia comparada con foco europeo y cluster intermedio OCDE",
    "Tres escenarios con trade-offs explícitos",
], size=14)
add_footer_brand(s)


# ----- Slide 4: Mensajes clave -----
s = prs.slides.add_slide(blank)
slide_title(s, "Mensajes clave")

# Cinco bullets numerados 1 a 5 (alineados con informe v4.3)
mensajes = [
    ("1", "Gasto de bolsillo estructural",
     "62 % del gasto total en medicamentos (CIF/UC, 2024) y 71 % del gasto retail (OECD SHA, 2022). La cifra más alta de OCDE en ambos cortes."),
    ("2", "Dos miradas complementarias",
     "Mirada acumulativa (crónicos ambulatorios). Mirada catastrófica (alto costo). Ninguna medida única atiende ambas a la vez."),
    ("3", "Patrón de paquete OCDE",
     "Los países que reducen sostenidamente el OOP integran al menos cuatro de seis componentes: canasta priorizada, copagos topados, dispensación accesible, sustitución por valor, regulación focal y trazabilidad."),
    ("4", "BFU como zoom analítico",
     "El informe profundiza el escenario de convergencia intermedia (E2). El Beneficio Farmacéutico Universal articula los regímenes existentes (GES, LRS, DAC, FOFAR, CAEC) bajo reglas comunes Fonasa-Isapre."),
    ("5", "Transparencia presupuestaria pendiente",
     "GES Fonasa medicamentos sin glosa segregada. CAEC porción farmacéutica no publicada. Brecha de información para diseñar política."),
]

y = Inches(1.2)
for num, title, body in mensajes:
    # Numero grande naranja
    add_text(s, Inches(0.6), y, Inches(0.6), Inches(1.0),
             num, 36, bold=True, color=COL_ACCENT, font=FONT_HEADING)
    add_text(s, Inches(1.3), y + Emu(50000), Inches(11.5), Inches(0.4),
             title, 15, bold=True, color=COL_PRIMARY)
    add_text(s, Inches(1.3), y + Inches(0.45), Inches(11.5), Inches(0.7),
             body, 11, color=COL_TEXT)
    y += Inches(1.15)

add_footer_brand(s)


# ----- Slide 5: Diagnóstico dual -----
s = prs.slides.add_slide(blank)
slide_title(s, "Diagnóstico: dos miradas, dos instrumentos")

# Dos columnas
add_rect(s, Inches(0.6), Inches(1.3), Inches(6), Inches(5.5), COL_BG_SOFT)
add_rect(s, Inches(6.8), Inches(1.3), Inches(6), Inches(5.5), COL_BG_SOFT)

add_text(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.6),
         "(1) Mirada acumulativa", 22, bold=True, color=COL_ACCENT, font=FONT_HEADING)
add_bullets(s, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.2), [
    "Crónicos ambulatorios: hipertensión, diabetes, salud mental, asma",
    "Gasto recurrente que se acumula mes a mes",
    "29 % dejó de tomar dosis por costo (Ipsos-EP, jul 2025)",
    "Ausencia de tope anual genera brecha estructural visible",
    "El precio unitario es solo una pieza; sin tope anual el costo se acumula",
], size=14)

add_text(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.6),
         "(2) Mirada catastrófica", 22, bold=True, color=COL_PRIMARY, font=FONT_HEADING)
add_bullets(s, Inches(7.0), Inches(2.2), Inches(5.6), Inches(4.2), [
    "Oncológicos modernos, enfermedades poco frecuentes, biotecnológicos",
    "Riesgo de agotar patrimonio familiar",
    "Cobertura fragmentada: GES, LRS, DAC, judicialización",
    "Ley Ricarte Soto: lista cerrada, ingresos discrecionales",
    "Cobertura universal nominal no elimina la judicialización (cf. Colombia)",
], size=14)

add_footer_brand(s)


# ----- Slide 6: Tabla de protección farmacéutica vigente -----
s = prs.slides.add_slide(blank)
slide_title(s, "Tabla de protección farmacéutica vigente")

instrumentos = [
    ("GES (medicamentos)", "Fonasa + Isapre", "87 problemas, sin glosa segregada (gap)"),
    ("Ley Ricarte Soto", "Fonasa + Isapre", "27 patologías de alto costo"),
    ("DAC Glosa 11", "Fonasa", "Oncológicas no GES/LRS"),
    ("FOFAR + Arsenal APS", "Fonasa APS", "HTA, DM2, dislipidemia, esenciales"),
    ("CAEC", "Isapre (seguro adicional)", "Catastrófico privado, porción farmacéutica sin publicar"),
    ("Cenabast Ley 21.198", "Hogares vía retail adherido", "Subsidio cruzado vía precio"),
]

# Header
y = Inches(1.3)
add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.45), COL_PRIMARY)
add_text(s, Inches(0.7), y + Emu(40000), Inches(3.5), Inches(0.35),
         "INSTRUMENTO", 12, bold=True, color=COL_WHITE, font=FONT_HEADING)
add_text(s, Inches(4.3), y + Emu(40000), Inches(3.0), Inches(0.35),
         "BENEFICIARIOS", 12, bold=True, color=COL_WHITE, font=FONT_HEADING)
add_text(s, Inches(7.4), y + Emu(40000), Inches(5.3), Inches(0.35),
         "COBERTURA", 12, bold=True, color=COL_WHITE, font=FONT_HEADING)
y += Inches(0.55)

for i, (inst, ben, cob) in enumerate(instrumentos):
    if i % 2 == 0:
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.7), COL_BG_SOFT)
    add_text(s, Inches(0.7), y + Inches(0.1), Inches(3.5), Inches(0.6),
             inst, 12, bold=True, color=COL_SECONDARY)
    add_text(s, Inches(4.3), y + Inches(0.1), Inches(3.0), Inches(0.6),
             ben, 11, color=COL_TEXT)
    add_text(s, Inches(7.4), y + Inches(0.1), Inches(5.3), Inches(0.6),
             cob, 11, color=COL_TEXT)
    y += Inches(0.78)

add_footer_brand(s)


# ----- Slide 7: Evidencia comparada -----
s = prs.slides.add_slide(blank)
slide_title(s, "Lo que muestra la evidencia comparada")
add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "Países agrupados por cluster fiscal de gasto público en medicamentos",
         15, italic=True, color=COL_MUTED)

# Tabla por cluster
y = Inches(1.85)
headers = ["Cluster", "Países", "Gasto público (% PIB)", "OOP (% gasto total)"]
widths_in = [3.4, 3.6, 2.8, 2.2]
xs = [Inches(0.6)]
for w in widths_in[:-1]:
    xs.append(xs[-1] + Inches(w))

add_rect(s, Inches(0.6), y, Inches(12), Inches(0.45), COL_PRIMARY)
for i, (h, w) in enumerate(zip(headers, widths_in)):
    add_text(s, xs[i] + Inches(0.05), y + Emu(40000), Inches(w), Inches(0.35),
             h, 12, bold=True, color=COL_WHITE, font=FONT_HEADING)
y += Inches(0.55)

clusters = [
    ("Alta cobertura", "Alemania, Francia, Países Bajos", "1,3 a 1,5 %", "menos de 25 %", False),
    ("Cluster intermedio OCDE", "España, Reino Unido, Canadá", "0,8 a 1,2 %", "26 a 43 %", False),
    ("Bajo esfuerzo fiscal", "Chile (actual)", "0,46 %", "62 a 71 %", True),
]
for i, row in enumerate(clusters):
    bg = COL_BG_SOFT if i % 2 == 0 else COL_WHITE
    add_rect(s, Inches(0.6), y, Inches(12), Inches(0.7), bg)
    is_chile = row[4]
    color = COL_ACCENT if is_chile else COL_TEXT
    bold = is_chile
    for j, w in enumerate(widths_in):
        val = row[j]
        add_text(s, xs[j] + Inches(0.05), y + Inches(0.18), Inches(w), Inches(0.5),
                 val, 12, bold=bold, color=color)
    y += Inches(0.78)

add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.6),
         "Chile invierte menos de la mitad del cluster intermedio OCDE; los hogares cubren la diferencia.",
         13, italic=True, bold=True, color=COL_PRIMARY)

add_footer_brand(s)


# ----- Slide 8: Tres escenarios -----
s = prs.slides.add_slide(blank)
slide_title(s, "Tres escenarios de política para Chile")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "Cifras en precios 2022 · órdenes de magnitud. El informe ordena la conversación.",
         14, italic=True, color=COL_MUTED)

escenarios = [
    {"name": "E1 · Ajuste gradual", "color": RGBColor(0x80, 0xA0, 0xC0),
     "rows": [("Lógica", "Profundizar instrumentos vigentes"),
              ("OOP esperado", "alrededor de 60 %"),
              ("Costo fiscal", "USD 400 a 500 M/año"),
              ("Reforma legal", "Modificación GES, LRS, DAC")]},
    {"name": "E2 · BFU intermedio", "color": RGBColor(0x80, 0x40, 0x60),
     "rows": [("Lógica", "Beneficio Farmacéutico Universal con cluster intermedio"),
              ("OOP esperado", "30 a 40 %"),
              ("Costo fiscal", "USD 800 a 900 M/año"),
              ("Reforma legal", "Ley marco BFU")]},
    {"name": "E3 · Convergencia plena", "color": RGBColor(0x60, 0x40, 0x80),
     "rows": [("Lógica", "Tope universal, paquete OCDE alto"),
              ("OOP esperado", "niveles bajos"),
              ("Costo fiscal", "USD 2.500 a 3.000 M/año"),
              ("Reforma legal", "Reforma sistémica integrada con Fonasa")]},
]

x = Inches(0.6)
col_w = Inches(4.05)
y_top = Inches(1.85)
for esc in escenarios:
    add_rect(s, x, y_top, col_w, Inches(0.5), esc["color"])
    add_text(s, x + Inches(0.15), y_top + Inches(0.1), col_w, Inches(0.4),
             esc["name"], 13, bold=True, color=COL_WHITE, font=FONT_HEADING)
    yy = y_top + Inches(0.55)
    for k, v in esc["rows"]:
        add_text(s, x + Inches(0.15), yy, col_w - Inches(0.3), Inches(0.3),
                 k.upper(), 9, bold=True, color=COL_MUTED, font=FONT_HEADING)
        add_text(s, x + Inches(0.15), yy + Inches(0.28), col_w - Inches(0.3), Inches(0.7),
                 v, 12, color=COL_TEXT)
        yy += Inches(0.95)
    x += col_w + Inches(0.15)

add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
         "El BFU (E2) se desarrolla con mayor detalle por riqueza analítica, no por preferencia.",
         12, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER)

add_footer_brand(s)


# ----- Slide 9: Zoom analítico BFU -----
s = prs.slides.add_slide(blank)
slide_title(s, "Zoom analítico: Beneficio Farmacéutico Universal")

add_text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
         "Componentes operacionales del Escenario 2", 16,
         italic=True, color=COL_MUTED)

# Dos columnas con componentes operacionales
add_rect(s, Inches(0.6), Inches(1.85), Inches(6.05), Inches(4.95), COL_BG_SOFT)
add_rect(s, Inches(6.85), Inches(1.85), Inches(6.05), Inches(4.95), COL_BG_SOFT)

add_text(s, Inches(0.8), Inches(2.0), Inches(5.8), Inches(0.5),
         "Diseño del beneficio", 16, bold=True, color=COL_PRIMARY, font=FONT_HEADING)
add_bullets(s, Inches(0.8), Inches(2.55), Inches(5.7), Inches(4.2), [
    "Lista positiva universal priorizada por ETESA (no negativa, no fragmentaria)",
    "Tope OOP del hogar calibrado en 13-15 % del ingreso per cápita (USD 800-900 M/año, simulación EPF anexo §5)",
    "Modalidad subsidio en POS (no reembolso); reduce barrera Q1-Q3",
    "Sustitución bioequivalente obligatoria con incentivo de precio",
], size=12)

add_text(s, Inches(7.05), Inches(2.0), Inches(5.8), Inches(0.5),
         "Articulación e implementación", 16, bold=True, color=COL_PRIMARY, font=FONT_HEADING)
add_bullets(s, Inches(7.05), Inches(2.55), Inches(5.7), Inches(4.2), [
    "Cobertura sobre todo OOP: retail y hospitalario ambulatorio",
    "Articulación con regímenes vigentes (GES, LRS, DAC, FOFAR, CAEC); el BFU es residual sobre lo que esos cubren",
    "Trazabilidad por RUT entre Fonasa e Isapre",
    "Acuerdos de acceso gestionado para innovación de alto costo",
], size=12)

add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
         "Capítulo 7 del informe v4.3 · simulación de costo fiscal en anexo §5.",
         12, italic=True, color=COL_PRIMARY, align=PP_ALIGN.CENTER)

add_footer_brand(s)


# ----- Slide 10: Triple-pack -----
s = prs.slides.add_slide(blank)
slide_title(s, "Lo que entrego en esta sesión")

add_rect(s, Inches(0.6), Inches(1.5), Inches(4.0), Inches(5.5), COL_BG_SOFT)
add_rect(s, Inches(4.7), Inches(1.5), Inches(4.0), Inches(5.5), COL_BG_SOFT)
add_rect(s, Inches(8.8), Inches(1.5), Inches(4.0), Inches(5.5), COL_BG_SOFT)

# Long doc
add_text(s, Inches(0.8), Inches(1.7), Inches(3.6), Inches(0.5),
         "INFORME EXTENSO", 14, bold=True, color=COL_PRIMARY, font=FONT_HEADING)
add_text(s, Inches(0.8), Inches(2.1), Inches(3.6), Inches(0.5),
         "v4.3", 32, bold=True, color=COL_SECONDARY, font=FONT_HEADING)
add_bullets(s, Inches(0.8), Inches(2.9), Inches(3.6), Inches(4.0), [
    "57 páginas",
    "8 capítulos + 6 anexos",
    "9 tablas, 33 figuras",
    "Marco BFU consolidado",
    "Cifras dobles 62 % / 71 %",
    "Tabla protección con 6 instrumentos",
    "Zoom BFU operacional",
    "Tracked changes preservados",
], size=12)

# Brief
add_text(s, Inches(4.9), Inches(1.7), Inches(3.6), Inches(0.5),
         "POLICY BRIEF", 14, bold=True, color=COL_PRIMARY, font=FONT_HEADING)
add_text(s, Inches(4.9), Inches(2.1), Inches(3.6), Inches(0.5),
         "v4.3", 32, bold=True, color=COL_SECONDARY, font=FONT_HEADING)
add_bullets(s, Inches(4.9), Inches(2.9), Inches(3.6), Inches(4.0), [
    "Extracto del informe",
    "Identidad visual EP",
    "Panorama y opciones",
    "5 mensajes clave",
    "Tres escenarios",
    "Zoom BFU",
    "Nota institucional concisa",
    "Para CIF y público",
], size=12)

# PPT
add_text(s, Inches(9.0), Inches(1.7), Inches(3.6), Inches(0.5),
         "ESTA PRESENTACIÓN", 14, bold=True, color=COL_PRIMARY, font=FONT_HEADING)
add_text(s, Inches(9.0), Inches(2.1), Inches(3.6), Inches(0.5),
         "v4 · 12 sl.", 32, bold=True, color=COL_SECONDARY, font=FONT_HEADING)
add_bullets(s, Inches(9.0), Inches(2.9), Inches(3.6), Inches(4.0), [
    "Para esta sesión deliberativa",
    "Ancla la discusión",
    "Resume el panorama",
    "Presenta opciones",
    "Espacio para feedback",
    "No es la presentación pública",
], size=12)

add_text(s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.35),
         "Informe v4.3 (57 pp) · Brief v4.3 (extracto del informe) · PPT v4 (12 sl)",
         11, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER)

add_footer_brand(s)


# ----- Slide 11: 5 puntos técnicos del mail -----
s = prs.slides.add_slide(blank)
slide_title(s, "Próximas decisiones")

add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
         "Cinco puntos técnicos sobre los que necesito su mirada", 22,
         bold=True, color=COL_PRIMARY, font=FONT_HEADING)

add_bullets(s, Inches(0.6), Inches(2.1), Inches(12), Inches(4.5), [
    "Fuente unificada para fichas país: Anexos 6/7 vs Tabla 4 del Capítulo 5",
    "Articulación operativa BFU + GES + LRS + DAC + FOFAR (Capítulo 7.3)",
    "Unidad del tope: persona u hogar (Capítulo 7.4.2)",
    "Clasificación de medidas por horizonte de reforma (Capítulo 8.1)",
    "CENABAST como compras agregadas dentro del beneficio (Ley 21.198)",
], size=15)

add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "Cinco puntos que definen la arquitectura del beneficio propuesto.",
         13, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER)

add_footer_brand(s)


# ----- Slide 12: Preguntas para el seminario -----
s = prs.slides.add_slide(blank)
slide_title(s, "Preguntas para la discusión pública")

add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
         "Cuatro dilemas que el informe abre, no cierra", 22,
         bold=True, color=COL_PRIMARY, font=FONT_HEADING, italic=True)

add_bullets(s, Inches(0.6), Inches(2.4), Inches(12), Inches(4.0), [
    "¿BFU sustituye o complementa los instrumentos vigentes?",
    "¿La unidad del tope debe ser persona u hogar?",
    "¿Qué medidas son de corto plazo y cuáles requieren reforma estructural?",
    "¿CENABAST debe ser parte del BFU o instrumento paralelo?",
], size=16)

add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "Gracias.",
         16, bold=True, italic=True, color=COL_PRIMARY, font=FONT_HEADING, align=PP_ALIGN.CENTER)

add_footer_brand(s)


# ----- Save -----
prs.save(str(OUT))
print(f"PPTX saved: {OUT}")
print(f"Size: {OUT.stat().st_size:,} bytes")
print(f"Slides: {len(prs.slides)}")
